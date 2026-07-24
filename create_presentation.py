"""
VOO Ward Admin Dashboard — Project Proposal Presentation
Professional proposal structure with formal tone.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Palette ────────────────────────────────────────────────────
PRIMARY      = RGBColor(0x1B, 0x2A, 0x4A)   # Navy blue
ACCENT       = RGBColor(0x2E, 0x86, 0xAB)   # Teal
ACCENT2      = RGBColor(0x00, 0x6D, 0x77)   # Dark teal
HIGHLIGHT    = RGBColor(0xE8, 0x6F, 0x51)   # Warm coral for emphasis
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF5, 0xF5, 0xF0)
LIGHT_GRAY   = RGBColor(0x8C, 0x8C, 0x96)
DARK_GRAY    = RGBColor(0x3A, 0x3A, 0x4A)
VERY_LIGHT   = RGBColor(0xE8, 0xE8, 0xEC)
GREEN_OK     = RGBColor(0x27, 0xAE, 0x60)
AMBER_WARN   = RGBColor(0xF3, 0x9C, 0x12)
RED_RISK     = RGBColor(0xC0, 0x39, 0x2B)
BG_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SECTION_BG   = RGBColor(0xF0, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper Functions ─────────────────────────────────────────────────

def fill_bg(slide, color=BG_WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def rect(slide, left, top, width, height, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def flat_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def txt(slide, left, top, width, height, text, size=18, color=DARK_GRAY,
        bold=False, align=PP_ALIGN.LEFT, font='Calibri', spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if spacing != 1.0:
        p.line_spacing = Pt(size * spacing)
    return box

def header_bar(slide, title, subtitle=None, section_num=None):
    fill_bg(slide)
    flat_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.35), PRIMARY)
    flat_rect(slide, Inches(0), Inches(1.35), Inches(13.333), Pt(4), ACCENT)
    if section_num:
        txt(slide, Inches(0.8), Inches(0.25), Inches(1), Inches(0.4),
            section_num, size=14, color=ACCENT, bold=True)
    txt(slide, Inches(0.8), Inches(0.45), Inches(10), Inches(0.6),
        title, size=32, color=WHITE, bold=True)
    if subtitle:
        txt(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.35),
            subtitle, size=14, color=RGBColor(0xBB, 0xCC, 0xDD))
    txt(slide, Inches(10.5), Inches(0.45), Inches(2.5), Inches(0.6),
        "VOO WARD", size=16, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)

def footer(slide):
    flat_rect(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), PRIMARY)
    txt(slide, Inches(0.8), Inches(7.18), Inches(5), Inches(0.3),
        "VOO Ward Admin Dashboard  |  Project Proposal", size=9, color=RGBColor(0x99, 0xAA, 0xBB))
    txt(slide, Inches(9), Inches(7.18), Inches(4), Inches(0.3),
        "Confidential", size=9, color=RGBColor(0x99, 0xAA, 0xBB), align=PP_ALIGN.RIGHT)

def card(slide, left, top, width, height, title, items, accent=ACCENT):
    rect(slide, left, top, width, height, SECTION_BG, border=VERY_LIGHT)
    flat_rect(slide, left, top, width, Pt(4), accent)
    txt(slide, left + Inches(0.25), top + Inches(0.18), width - Inches(0.5), Inches(0.35),
        title, size=14, color=PRIMARY, bold=True)
    y = top + Inches(0.55)
    for item in items:
        txt(slide, left + Inches(0.35), y, width - Inches(0.6), Inches(0.26),
            item, size=11, color=DARK_GRAY)
        y += Inches(0.26)


# =====================================================================
# SLIDE 1 — TITLE / COVER
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, PRIMARY)
flat_rect(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5), ACCENT)

txt(slide, Inches(1.5), Inches(1.2), Inches(6), Inches(0.5),
    "PROJECT PROPOSAL", size=16, color=ACCENT, bold=True)
flat_rect(slide, Inches(1.5), Inches(1.75), Inches(3.5), Pt(3), ACCENT)

txt(slide, Inches(1.5), Inches(2.1), Inches(8), Inches(1.2),
    "VOO Ward\nAdmin Dashboard", size=52, color=WHITE, bold=True)

txt(slide, Inches(1.5), Inches(4.0), Inches(7), Inches(0.8),
    "A Comprehensive Digital Governance Platform\nfor Kyamatu Ward Administration",
    size=20, color=RGBColor(0xBB, 0xCC, 0xDD))

flat_rect(slide, Inches(1.5), Inches(5.3), Inches(5), Pt(1), RGBColor(0x3A, 0x5A, 0x7A))
meta = [
    ("Prepared for:", "Kyamatu Ward, County Government"),
    ("Prepared by:", "VOO Ward Development Team"),
    ("Date:", "February 2026"),
    ("Document:", "Technical Proposal v1.0"),
]
y = Inches(5.5)
for label, value in meta:
    txt(slide, Inches(1.5), y, Inches(1.8), Inches(0.28),
        label, size=11, color=LIGHT_GRAY)
    txt(slide, Inches(3.3), y, Inches(4), Inches(0.28),
        value, size=11, color=WHITE, bold=True)
    y += Inches(0.3)

# Right side — executive summary
rect(slide, Inches(8.5), Inches(1.5), Inches(4.2), Inches(5.2),
     RGBColor(0x22, 0x3A, 0x5A), border=RGBColor(0x3A, 0x5A, 0x7A))
flat_rect(slide, Inches(8.5), Inches(1.5), Inches(4.2), Pt(4), ACCENT)
txt(slide, Inches(8.8), Inches(1.75), Inches(3.6), Inches(0.3),
    "EXECUTIVE SUMMARY", size=12, color=ACCENT, bold=True)
txt(slide, Inches(8.8), Inches(2.2), Inches(3.6), Inches(4.2),
    "This proposal presents a web-based Admin Dashboard for Kyamatu Ward that enables the Member of County Assembly (MCA) and administrative staff to:\n\n"
    "- Track and resolve citizen-reported issues\n"
    "- Process student bursary applications\n"
    "- Manage constituent registration data\n"
    "- Broadcast announcements via web and USSD\n"
    "- Monitor ward operations through real-time analytics\n"
    "- Maintain full audit trails of all actions\n\n"
    "The platform integrates with USSD for feature phone access and a mobile app for smartphone users.",
    size=11, color=RGBColor(0xCC, 0xDD, 0xEE), spacing=1.5)


# =====================================================================
# SLIDE 2 — TABLE OF CONTENTS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Table of Contents", "Proposal Structure")

sections = [
    ("01", "Background and Problem Statement", "Context and challenges facing ward administration"),
    ("02", "Project Objectives", "Goals and measurable outcomes"),
    ("03", "Proposed Solution", "System overview and key modules"),
    ("04", "Functional Requirements", "Detailed feature specifications"),
    ("05", "System Architecture", "Technical design and technology stack"),
    ("06", "Security Framework", "Authentication, access control and data protection"),
    ("07", "User Roles and Access Control", "Role-based permissions matrix"),
    ("08", "Analytics and Reporting", "Data visualization and export capabilities"),
    ("09", "Implementation Approach", "Methodology, phases and timeline"),
    ("10", "Expected Outcomes and Benefits", "Measurable impact and value delivered"),
    ("11", "Risk Assessment and Mitigation", "Identified risks and contingency plans"),
    ("12", "Conclusion and Recommendations", "Summary and next steps"),
]

col1_x, col2_x = Inches(0.8), Inches(6.8)
y_start = Inches(1.8)
for i, (num, title, desc) in enumerate(sections):
    x = col1_x if i < 6 else col2_x
    y = y_start + (i % 6) * Inches(0.85)
    rect(slide, x, y, Inches(0.42), Inches(0.42), PRIMARY if i < 6 else ACCENT)
    txt(slide, x + Inches(0.02), y + Pt(2), Inches(0.4), Inches(0.38),
        num, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.55), y, Inches(5), Inches(0.28),
        title, size=14, color=PRIMARY, bold=True)
    txt(slide, x + Inches(0.55), y + Inches(0.3), Inches(5), Inches(0.28),
        desc, size=11, color=LIGHT_GRAY)
footer(slide)


# =====================================================================
# SLIDE 3 — BACKGROUND AND PROBLEM STATEMENT
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Background and Problem Statement", "Understanding the challenges in ward administration", "01")

txt(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.9),
    "Kyamatu Ward, like many administrative units in Kenya, faces significant challenges in managing citizen services, tracking public issues, and ensuring transparent allocation of bursary funds. Current manual and fragmented processes lead to delays, lack of accountability, and limited citizen engagement.",
    size=14, color=DARK_GRAY, spacing=1.6)

problems = [
    ("Fragmented Issue Tracking",
     "Citizen complaints are received via phone calls, in-person visits, and informal channels with no centralized system for logging, tracking, or resolving them. Issues fall through the cracks and citizens have no visibility into resolution status."),
    ("Opaque Bursary Processing",
     "Bursary applications are handled manually with paper-based records. There is no audit trail for approvals or rejections, making the process vulnerable to errors, delays, and lack of transparency."),
    ("Limited Citizen Communication",
     "Ward announcements rely on word-of-mouth and notice boards. Citizens with feature phones (the majority) have no digital channel to receive updates, report issues, or check service status."),
    ("No Data-Driven Decision Making",
     "Without centralized data, the MCA and staff cannot analyze issue trends, measure resolution rates, assess bursary distribution patterns, or identify priority areas needing attention."),
]

y = Inches(2.8)
for i, (title, desc) in enumerate(problems):
    x = Inches(0.8) if i % 2 == 0 else Inches(6.8)
    if i == 2:
        y = Inches(5.0)
    rect(slide, x, y, Inches(5.7), Inches(1.8), SECTION_BG, border=VERY_LIGHT)
    flat_rect(slide, x, y, Pt(5), Inches(1.8), HIGHLIGHT if i < 2 else AMBER_WARN)
    txt(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.1), Inches(0.3),
        title, size=14, color=PRIMARY, bold=True)
    txt(slide, x + Inches(0.3), y + Inches(0.5), Inches(5.1), Inches(1.1),
        desc, size=11, color=DARK_GRAY, spacing=1.5)
footer(slide)


# =====================================================================
# SLIDE 4 — PROJECT OBJECTIVES
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Project Objectives", "Goals and measurable outcomes", "02")

objectives = [
    ("Centralize Issue Management",
     "Provide a single platform to log, categorize, assign, track, and resolve citizen-reported issues from all channels (web, USSD, mobile) with full status visibility."),
    ("Digitize Bursary Processing",
     "Implement a transparent, auditable workflow for receiving, reviewing, approving/rejecting, and disbursing student bursary applications with role-based access controls."),
    ("Enable Multi-Channel Communication",
     "Allow the ward office to broadcast announcements through the web dashboard, USSD service, SMS, and mobile app, reaching citizens regardless of their device type."),
    ("Establish Data-Driven Governance",
     "Deliver real-time analytics dashboards with charts, maps, and exportable reports to support evidence-based decision making and performance monitoring."),
    ("Ensure Security and Accountability",
     "Implement role-based access control, comprehensive audit trails, secure authentication, and data protection measures meeting best-practice standards."),
    ("Maintain a Constituent Registry",
     "Build and manage a verified database of ward constituents with registration, identity verification, and location-based grouping capabilities."),
]

y = Inches(1.7)
col1_x, col2_x = Inches(0.8), Inches(6.8)
for i, (title, desc) in enumerate(objectives):
    x = col1_x if i % 2 == 0 else col2_x
    if i > 0 and i % 2 == 0:
        y += Inches(1.65)
    rect(slide, x, y, Inches(5.7), Inches(1.5), SECTION_BG, border=VERY_LIGHT)
    rect(slide, x + Inches(0.2), y + Inches(0.18), Inches(0.38), Inches(0.38), ACCENT)
    txt(slide, x + Inches(0.22), y + Inches(0.2), Inches(0.36), Inches(0.35),
        str(i + 1), size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.7), y + Inches(0.15), Inches(4.7), Inches(0.3),
        title, size=13, color=PRIMARY, bold=True)
    txt(slide, x + Inches(0.7), y + Inches(0.48), Inches(4.7), Inches(0.9),
        desc, size=11, color=DARK_GRAY, spacing=1.45)
footer(slide)


# =====================================================================
# SLIDE 5 — PROPOSED SOLUTION OVERVIEW
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Proposed Solution", "System overview and core modules", "03")

txt(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.6),
    "The VOO Ward Admin Dashboard is a web-based platform that serves as the central hub for all ward administrative operations. It integrates with a USSD service for feature phone users and a React Native mobile application for smartphone users.",
    size=13, color=DARK_GRAY, spacing=1.5)

modules = [
    ("Issue Management", ACCENT, [
        "Centralized ticket system with unique IDs",
        "Multi-source intake (web, USSD, mobile)",
        "Status workflow: Open > In Progress > Resolved > Closed",
        "Priority classification (Urgent/High/Medium/Low)",
        "Bulk resolution capability for MCA",
    ]),
    ("Bursary Processing", GREEN_OK, [
        "End-to-end application lifecycle",
        "Review and approval workflow",
        "Amount tracking in KES",
        "Decision documentation with notes",
        "MCA-restricted access control",
    ]),
    ("Announcements", AMBER_WARN, [
        "Multi-type: General, Urgent, Event, Service Update",
        "USSD news feed integration",
        "Optional SMS broadcast on publish",
        "Citizen message inbox with reply",
        "Timestamped publication history",
    ]),
    ("Constituent Registry", ACCENT2, [
        "Multi-channel voter registration",
        "Identity verification workflow",
        "National ID and location tracking",
        "Village and sub-location grouping",
        "Exportable constituent database",
    ]),
    ("Analytics and Reporting", PRIMARY, [
        "Interactive charts (Chart.js)",
        "Geographic issue mapping (Leaflet.js)",
        "Date range filtered reports",
        "CSV export with Excel compatibility",
        "Resolution rate and trend analysis",
    ]),
    ("User Administration", HIGHLIGHT, [
        "Three-tier role system (MCA/PA/Clerk)",
        "Self-service signup for staff",
        "Profile management with avatar upload",
        "Session management and security",
        "Comprehensive audit trail",
    ]),
]

x_start = Inches(0.8)
y_top = Inches(2.6)
col_w = Inches(3.75)
gap = Inches(0.2)
for i, (title, color, items) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = x_start + col * (col_w + gap)
    y = y_top + row * Inches(2.55)
    rect(slide, x, y, col_w, Inches(2.35), SECTION_BG, border=VERY_LIGHT)
    flat_rect(slide, x, y, col_w, Pt(4), color)
    txt(slide, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.4), Inches(0.3),
        title, size=13, color=PRIMARY, bold=True)
    iy = y + Inches(0.5)
    for item in items:
        txt(slide, x + Inches(0.3), iy, col_w - Inches(0.5), Inches(0.24),
            "- " + item, size=10, color=DARK_GRAY)
        iy += Inches(0.24)
footer(slide)


# =====================================================================
# SLIDE 6 — FUNCTIONAL REQUIREMENTS (ISSUE MANAGEMENT)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Functional Requirements", "Issue Management — Detailed Specifications", "04a")

reqs = [
    ("FR-01", "Issue Creation", "Admin can create issues via modal form with title, description, priority level, category, and location fields."),
    ("FR-02", "Multi-Source Intake", "Issues can originate from the web dashboard, USSD service, or mobile application. Source is automatically tagged."),
    ("FR-03", "Status Management", "Issues follow a defined workflow: Open > In Progress > Resolved > Closed. Each transition requires action notes."),
    ("FR-04", "Bulk Resolution", "MCA can select and resolve up to 200 issues simultaneously with a single action and shared resolution note."),
    ("FR-05", "Category Classification", "Predefined categories include Infrastructure, Water, Health, and custom categories. Used for filtering and analytics."),
    ("FR-06", "Priority Levels", "Four priority levels (Urgent, High, Medium, Low) determine visibility and urgency in the dashboard and map."),
    ("FR-07", "Search and Filter", "Filter issues by status, category, priority, date range, and source. Full-text search across issue descriptions."),
    ("FR-08", "CSV Export", "Export filtered issue data to CSV format with UTF-8 BOM encoding for Microsoft Excel compatibility."),
    ("FR-09", "USSD Synchronization", "Issues reported via USSD are automatically reflected in the dashboard. Resolution updates propagate to USSD."),
    ("FR-10", "Status Summary", "Dashboard footer displays real-time counts of pending, in-progress, and resolved issues."),
]

flat_rect(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.4), PRIMARY)
txt(slide, Inches(0.9), Inches(1.68), Inches(1.2), Inches(0.35), "ID", size=11, color=WHITE, bold=True)
txt(slide, Inches(2.1), Inches(1.68), Inches(2.5), Inches(0.35), "Requirement", size=11, color=WHITE, bold=True)
txt(slide, Inches(4.6), Inches(1.68), Inches(7.5), Inches(0.35), "Description", size=11, color=WHITE, bold=True)

y = Inches(2.1)
for i, (req_id, req_name, req_desc) in enumerate(reqs):
    bg = SECTION_BG if i % 2 == 0 else BG_WHITE
    flat_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.47), bg)
    txt(slide, Inches(0.9), y + Pt(3), Inches(1.2), Inches(0.4), req_id, size=10, color=ACCENT, bold=True, font='Consolas')
    txt(slide, Inches(2.1), y + Pt(3), Inches(2.5), Inches(0.4), req_name, size=10, color=PRIMARY, bold=True)
    txt(slide, Inches(4.6), y + Pt(3), Inches(7.5), Inches(0.4), req_desc, size=10, color=DARK_GRAY)
    y += Inches(0.47)
footer(slide)


# =====================================================================
# SLIDE 7 — FUNCTIONAL REQUIREMENTS (BURSARY & ANNOUNCEMENTS)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Functional Requirements", "Bursary Processing and Ward Communication", "04b")

txt(slide, Inches(0.8), Inches(1.65), Inches(5), Inches(0.3),
    "Bursary Application Management", size=16, color=PRIMARY, bold=True)
flat_rect(slide, Inches(0.8), Inches(1.97), Inches(2), Pt(3), ACCENT)

bursary_reqs = [
    ("FR-11", "Application Intake", "Process applications with applicant name, ID, institution, and requested amount in KES."),
    ("FR-12", "Review Workflow", "Applications follow: Pending > Under Review > Approved/Rejected > Disbursed."),
    ("FR-13", "Decision Documentation", "Every approval or rejection requires notes documenting the rationale."),
    ("FR-14", "MCA-Only Access", "Bursary viewing, processing, and export restricted to MCA role only."),
    ("FR-15", "Bursary Analytics", "Bar chart showing application distribution by status. Approval rate tracking."),
]

y = Inches(2.2)
for req_id, req_name, req_desc in bursary_reqs:
    txt(slide, Inches(0.9), y, Inches(1.2), Inches(0.28), req_id, size=10, color=ACCENT, bold=True, font='Consolas')
    txt(slide, Inches(2.1), y, Inches(2), Inches(0.28), req_name, size=10, color=PRIMARY, bold=True)
    txt(slide, Inches(4.2), y, Inches(8), Inches(0.28), req_desc, size=10, color=DARK_GRAY)
    y += Inches(0.35)

txt(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.3),
    "Announcements and Citizen Communication", size=16, color=PRIMARY, bold=True)
flat_rect(slide, Inches(0.8), Inches(4.62), Inches(2), Pt(3), ACCENT)

announcement_reqs = [
    ("FR-16", "Announcement Types", "Support four types: General Information, Urgent Notice, Event Notification, Service Update."),
    ("FR-17", "USSD News Feed", "Announcements automatically formatted and paginated for USSD screen display."),
    ("FR-18", "SMS Broadcast", "Optional SMS notification to registered constituents on announcement publish."),
    ("FR-19", "Citizen Inbox", "Public-facing message submission (no auth). Admin views, reads, and replies to messages."),
    ("FR-20", "USSD Interaction Logs", "All USSD sessions are logged with phone number, input, response, IP, and timestamp."),
]

y = Inches(4.85)
for req_id, req_name, req_desc in announcement_reqs:
    txt(slide, Inches(0.9), y, Inches(1.2), Inches(0.28), req_id, size=10, color=ACCENT, bold=True, font='Consolas')
    txt(slide, Inches(2.1), y, Inches(2), Inches(0.28), req_name, size=10, color=PRIMARY, bold=True)
    txt(slide, Inches(4.2), y, Inches(8), Inches(0.28), req_desc, size=10, color=DARK_GRAY)
    y += Inches(0.35)
footer(slide)


# =====================================================================
# SLIDE 8 — SYSTEM ARCHITECTURE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "System Architecture", "Technical design and technology stack", "05")

# Architecture tiers
rect(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(1.8), SECTION_BG, border=VERY_LIGHT)
txt(slide, Inches(0.8), Inches(1.75), Inches(3.5), Inches(0.25),
    "PRESENTATION TIER", size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
for j, c in enumerate(["Web Admin Dashboard (HTML/CSS/JS)", "USSD Service (*384#)", "Mobile App (React Native)"]):
    rect(slide, Inches(1.1), Inches(2.15) + j * Inches(0.48), Inches(2.9), Inches(0.4), PRIMARY)
    txt(slide, Inches(1.2), Inches(2.18) + j * Inches(0.48), Inches(2.7), Inches(0.35),
        c, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(slide, Inches(4.5), Inches(2.6), Inches(0.8), Inches(0.5),
    ">>>", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

rect(slide, Inches(5.0), Inches(2.0), Inches(3.5), Inches(1.8), SECTION_BG, border=VERY_LIGHT)
txt(slide, Inches(5.0), Inches(1.75), Inches(3.5), Inches(0.25),
    "APPLICATION TIER", size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
for j, a in enumerate(["Node.js + Express.js API", "Authentication Middleware", "Business Logic Layer"]):
    rect(slide, Inches(5.3), Inches(2.15) + j * Inches(0.48), Inches(2.9), Inches(0.4), ACCENT2)
    txt(slide, Inches(5.4), Inches(2.18) + j * Inches(0.48), Inches(2.7), Inches(0.35),
        a, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(slide, Inches(8.7), Inches(2.6), Inches(0.8), Inches(0.5),
    ">>>", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

rect(slide, Inches(9.2), Inches(2.0), Inches(3.5), Inches(1.8), SECTION_BG, border=VERY_LIGHT)
txt(slide, Inches(9.2), Inches(1.75), Inches(3.5), Inches(0.25),
    "DATA TIER", size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
for j, (d, c) in enumerate([("PostgreSQL (Supabase)", GREEN_OK), ("MongoDB Atlas", AMBER_WARN), ("Redis Cache", RED_RISK)]):
    rect(slide, Inches(9.5), Inches(2.15) + j * Inches(0.48), Inches(2.9), Inches(0.4), c)
    txt(slide, Inches(9.6), Inches(2.18) + j * Inches(0.48), Inches(2.7), Inches(0.35),
        d, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Tech stack table
txt(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.3),
    "Technology Stack Summary", size=16, color=PRIMARY, bold=True)
flat_rect(slide, Inches(0.8), Inches(4.55), Inches(2), Pt(3), ACCENT)

stack = [
    ("Runtime", "Node.js"), ("Framework", "Express.js"),
    ("Primary Database", "PostgreSQL via Supabase (managed)"),
    ("Secondary Database", "MongoDB Atlas (managed)"),
    ("Caching", "Redis"), ("Authentication", "bcryptjs + Bearer token sessions"),
    ("Charts", "Chart.js"), ("Maps", "Leaflet.js with OpenStreetMap"),
    ("Image Processing", "Sharp (server-side resize and format)"),
    ("File Storage", "AWS S3 v3 SDK (optional)"),
    ("Deployment", "Docker on Render.com"),
    ("Frontend", "Vanilla HTML, CSS, JavaScript (no framework)"),
]

flat_rect(slide, Inches(0.8), Inches(4.75), Inches(11.5), Inches(0.35), PRIMARY)
txt(slide, Inches(0.9), Inches(4.78), Inches(3), Inches(0.3), "Component", size=11, color=WHITE, bold=True)
txt(slide, Inches(4.2), Inches(4.78), Inches(8), Inches(0.3), "Technology", size=11, color=WHITE, bold=True)

y = Inches(5.15)
for i, (comp, tech) in enumerate(stack):
    bg = SECTION_BG if i % 2 == 0 else BG_WHITE
    flat_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.3), bg)
    txt(slide, Inches(0.9), y + Pt(1), Inches(3), Inches(0.28), comp, size=10, color=PRIMARY, bold=True)
    txt(slide, Inches(4.2), y + Pt(1), Inches(8), Inches(0.28), tech, size=10, color=DARK_GRAY)
    y += Inches(0.3)
footer(slide)


# =====================================================================
# SLIDE 9 — SECURITY FRAMEWORK
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Security Framework", "Authentication, data protection and compliance measures", "06")

card(slide, Inches(0.8), Inches(1.7), Inches(3.75), Inches(2.5),
     "Authentication", [
         "- Username and password-based login",
         "- bcrypt password hashing (10 salt rounds)",
         "- Auto-migration from legacy SHA-256 hashes",
         "- Bearer token session management",
         "- 30-minute server-side session timeout",
         "- 24-hour client-side token expiry",
         "- Single active session per user enforced",
     ], PRIMARY)

card(slide, Inches(4.75), Inches(1.7), Inches(3.75), Inches(2.5),
     "Rate Limiting and Brute Force Protection", [
         "- Server: 10 login attempts per minute per IP",
         "- Client: 5 failed attempts triggers 5-min lockout",
         "- express-rate-limit with in-memory fallback",
         "- Automatic session cleanup every 15 minutes",
         "- Login attempt logging for audit purposes",
         "", "",
     ], HIGHLIGHT)

card(slide, Inches(8.7), Inches(1.7), Inches(3.75), Inches(2.5),
     "Input Validation and Sanitization", [
         "- Server-side and client-side validation",
         "- HTML tag stripping on all text inputs",
         "- JavaScript protocol injection prevention",
         "- Event handler attribute removal",
         "- Email, phone, and username format checks",
         "- Request body size limit (10MB)",
         "",
     ], ACCENT)

card(slide, Inches(0.8), Inches(4.5), Inches(5.7), Inches(2.3),
     "HTTP Security Headers", [
         "- X-Frame-Options: DENY (clickjacking prevention)",
         "- X-Content-Type-Options: nosniff (MIME sniffing prevention)",
         "- X-XSS-Protection: 1; mode=block",
         "- Referrer-Policy: strict-origin-when-cross-origin",
         "- Content-Security-Policy with whitelisted sources",
         "- Cache-Control: no-cache, no-store, must-revalidate",
         "- CORS restricted to whitelisted origins only",
     ], ACCENT2)

card(slide, Inches(6.7), Inches(4.5), Inches(5.8), Inches(2.3),
     "Audit Trail and Accountability", [
         "- All admin actions logged: CREATE, UPDATE, DELETE, LOGIN, LOGOUT, EXPORT",
         "- Each log entry records: user, role, IP address, timestamp, description",
         "- Timeline view with filtering by action, user, and date range",
         "- Paginated audit log view (20 entries per page)",
         "- Audit logs exportable to CSV for external review",
         "- Protected MCA account (cannot be deleted from system)",
         "",
     ], PRIMARY)
footer(slide)


# =====================================================================
# SLIDE 10 — USER ROLES AND ACCESS CONTROL
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "User Roles and Access Control", "Role-based permissions matrix", "07")

txt(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.4),
    "The system enforces a three-tier role hierarchy with a maximum of three concurrent admin users: one MCA (super admin) and up to two PA or Clerk accounts.",
    size=13, color=DARK_GRAY, spacing=1.5)

flat_rect(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(0.45), PRIMARY)
for x, h in [(Inches(0.9), "Function"), (Inches(5.5), "MCA (Full Access)"), (Inches(8.0), "PA (Personal Assistant)"), (Inches(10.5), "Clerk")]:
    txt(slide, x, Inches(2.33), Inches(2.5), Inches(0.4), h, size=11, color=WHITE, bold=True)

permissions = [
    ("View Dashboard and Statistics", "Yes", "Yes", "Yes"),
    ("Create and Manage Issues", "Yes", "Yes", "Yes"),
    ("Create Announcements", "Yes", "Yes", "Yes"),
    ("Bulk Issue Resolution (up to 200)", "Yes", "No", "No"),
    ("View and Process Bursary Applications", "Yes", "No", "No"),
    ("Export Issues and USSD Data (CSV)", "Yes", "Yes", "No"),
    ("Export Bursaries and Constituents (CSV)", "Yes", "No", "No"),
    ("Manage Admin Users (Add/Edit/Delete)", "Yes", "No", "No"),
    ("Edit AI Chatbot Knowledge Base", "Yes", "No", "No"),
    ("View Audit Trail", "Yes", "Yes", "No"),
    ("Manage Own Profile and Password", "Yes", "Yes", "Yes"),
    ("Verify or Reject Constituents", "Yes", "Yes", "No"),
]

y = Inches(2.8)
for i, (func, mca, pa, clerk) in enumerate(permissions):
    bg = SECTION_BG if i % 2 == 0 else BG_WHITE
    flat_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.34), bg)
    txt(slide, Inches(0.9), y + Pt(1), Inches(4.5), Inches(0.3), func, size=10, color=DARK_GRAY)
    for val, vx in [(mca, Inches(5.5)), (pa, Inches(8.0)), (clerk, Inches(10.5))]:
        clr = GREEN_OK if val == "Yes" else RED_RISK
        txt(slide, vx, y + Pt(1), Inches(2), Inches(0.3), val, size=10, color=clr, bold=True)
    y += Inches(0.34)

rect(slide, Inches(0.8), Inches(6.95), Inches(11.5), Inches(0.5), SECTION_BG, border=VERY_LIGHT)
txt(slide, Inches(1.0), Inches(6.98), Inches(11), Inches(0.4),
    "Constraint: Maximum 3 admin users total (1 MCA + 2 PA/Clerk). The MCA account is protected and cannot be deleted. New PA/Clerk accounts are created through a self-service signup page.",
    size=10, color=DARK_GRAY)
footer(slide)


# =====================================================================
# SLIDE 11 — ANALYTICS AND REPORTING
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Analytics and Reporting", "Data visualization, metrics and export capabilities", "08")

card(slide, Inches(0.8), Inches(1.7), Inches(3.75), Inches(2.3),
     "Issue Status Distribution", [
         "- Doughnut chart (60% cutout)",
         "- Segments: Pending, In Progress, Resolved, Closed",
         "- Color-coded with interactive tooltips",
         "- Responsive to date range filters",
         "- Powered by Chart.js library",
     ], PRIMARY)

card(slide, Inches(4.75), Inches(1.7), Inches(3.75), Inches(2.3),
     "Issues Over Time", [
         "- Dual-dataset line chart",
         "- Tracks new issues vs. resolved issues",
         "- 30-day rolling window display",
         "- Trend analysis for planning",
         "- Responsive canvas rendering",
     ], ACCENT)

card(slide, Inches(8.7), Inches(1.7), Inches(3.75), Inches(2.3),
     "Bursary Applications Summary", [
         "- Bar chart by application status",
         "- Pending, Approved, Rejected, Received",
         "- Amount tracking per category",
         "- Approval rate percentage display",
         "- MCA-only access enforced",
     ], GREEN_OK)

card(slide, Inches(0.8), Inches(4.3), Inches(5.7), Inches(2.5),
     "Geographic Issue Mapping (Leaflet.js)", [
         "- OpenStreetMap tiles centered on Kyamatu Ward",
         "- Color-coded markers by issue status",
         "- Priority-level marker indicators",
         "- Interactive popups with issue details",
         "- Map statistics panel (total, pending, in-progress, resolved)",
         "- Filterable by status and priority level",
         "- Click-through to issue detail view",
     ], AMBER_WARN)

card(slide, Inches(6.7), Inches(4.3), Inches(5.8), Inches(2.5),
     "Data Export Capabilities", [
         "- Issues: ticket, category, status, message, phone, source, dates",
         "- Bursaries: ref code, student, institution, amount, status (MCA only)",
         "- Constituents: phone, ID, name, location, village, status (MCA only)",
         "- USSD Interactions: phone, input, response, ref code, IP, timestamp",
         "- Audit Logs: action type, user, role, IP, timestamp, description",
         "- All exports: CSV format with UTF-8 BOM for Excel compatibility",
         "- Date range filtering available across all export types",
     ], PRIMARY)
footer(slide)


# =====================================================================
# SLIDE 12 — IMPLEMENTATION APPROACH
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Implementation Approach", "Development methodology, phases and timeline", "09")

txt(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.5),
    "The project follows an Agile development methodology with iterative delivery across four phases. Each phase produces deployable increments reviewed with stakeholders before proceeding.",
    size=13, color=DARK_GRAY, spacing=1.5)

phases = [
    ("Phase 1", "Foundation", "Weeks 1-3", PRIMARY, [
        "Database schema design and migration setup",
        "Express.js API with authentication middleware",
        "Admin login, session management, role system",
        "Basic dashboard layout with sidebar navigation",
    ]),
    ("Phase 2", "Core Modules", "Weeks 4-7", ACCENT, [
        "Issue management (CRUD, status workflow, filters)",
        "Bursary processing (intake, review, approval)",
        "Announcement system and USSD news feed",
        "Constituent registry and verification",
    ]),
    ("Phase 3", "Analytics and Integration", "Weeks 8-10", ACCENT2, [
        "Chart.js analytics (doughnut, line, bar charts)",
        "Leaflet.js interactive issue map",
        "USSD service integration and logging",
        "CSV export and reporting capabilities",
    ]),
    ("Phase 4", "Polish and Deployment", "Weeks 11-12", GREEN_OK, [
        "AI chat assistant (Mai) integration",
        "Real-time notifications and polling",
        "Security hardening and audit trail",
        "Docker containerization and Render deployment",
    ]),
]

# Timeline bar
flat_rect(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.5), SECTION_BG)
phase_width = Inches(11.5) / 4
for i, (phase, name, weeks, color, _) in enumerate(phases):
    x = Inches(0.8) + phase_width * i
    flat_rect(slide, x + Pt(2), Inches(2.52), phase_width - Pt(4), Inches(0.46), color)
    txt(slide, x + Pt(4), Inches(2.52), phase_width - Pt(8), Inches(0.22),
        f"{phase}: {name}", size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Pt(4), Inches(2.76), phase_width - Pt(8), Inches(0.2),
        weeks, size=9, color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

y = Inches(3.3)
for i, (phase, name, weeks, color, items) in enumerate(phases):
    x = Inches(0.8) if i % 2 == 0 else Inches(6.8)
    if i == 2:
        y = Inches(5.3)
    rect(slide, x, y, Inches(5.7), Inches(1.7), SECTION_BG, border=VERY_LIGHT)
    flat_rect(slide, x, y, Inches(5.7), Pt(4), color)
    txt(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.2), Inches(0.28),
        f"{phase}: {name} ({weeks})", size=12, color=PRIMARY, bold=True)
    iy = y + Inches(0.48)
    for item in items:
        txt(slide, x + Inches(0.3), iy, Inches(5.1), Inches(0.24),
            "- " + item, size=10, color=DARK_GRAY)
        iy += Inches(0.26)
footer(slide)


# =====================================================================
# SLIDE 13 — EXPECTED OUTCOMES AND BENEFITS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Expected Outcomes and Benefits", "Measurable impact and value delivered", "10")

outcomes = [
    ("Improved Response Time",
     "Centralized issue tracking with priority classification enables the ward office to respond to urgent citizen concerns within hours rather than days or weeks.",
     "Target: 80% of urgent issues acknowledged within 4 hours"),
    ("Transparent Bursary Allocation",
     "Digital processing with audit trails ensures every bursary decision is documented, reducing disputes and increasing public trust in the allocation process.",
     "Target: 100% of decisions documented with rationale"),
    ("Expanded Citizen Reach",
     "USSD integration ensures citizens with feature phones (estimated 60%+ of ward residents) can report issues, receive announcements, and check service status.",
     "Target: 3x increase in citizen engagement within 6 months"),
    ("Data-Driven Decision Making",
     "Real-time analytics and geographic mapping allow the MCA to identify priority areas, allocate resources effectively, and measure resolution performance.",
     "Target: Monthly data-driven reports to County Assembly"),
    ("Operational Efficiency",
     "Automated workflows, bulk operations, and CSV exports reduce manual administrative work, freeing staff time for direct citizen engagement.",
     "Target: 50% reduction in administrative processing time"),
    ("Accountability and Compliance",
     "Comprehensive audit trails, role-based access, and security measures provide a verifiable record of all administrative actions taken in the system.",
     "Target: Full audit compliance for external review"),
]

y = Inches(1.7)
for i, (title, desc, metric) in enumerate(outcomes):
    x = Inches(0.8) if i % 2 == 0 else Inches(6.8)
    if i > 0 and i % 2 == 0:
        y += Inches(1.75)
    rect(slide, x, y, Inches(5.7), Inches(1.6), SECTION_BG, border=VERY_LIGHT)
    flat_rect(slide, x, y, Pt(5), Inches(1.6), ACCENT if i % 2 == 0 else PRIMARY)
    txt(slide, x + Inches(0.3), y + Inches(0.1), Inches(5.1), Inches(0.28),
        title, size=13, color=PRIMARY, bold=True)
    txt(slide, x + Inches(0.3), y + Inches(0.4), Inches(5.1), Inches(0.7),
        desc, size=10, color=DARK_GRAY, spacing=1.45)
    txt(slide, x + Inches(0.3), y + Inches(1.15), Inches(5.1), Inches(0.28),
        metric, size=10, color=ACCENT, bold=True)
footer(slide)


# =====================================================================
# SLIDE 14 — RISK ASSESSMENT AND MITIGATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Risk Assessment and Mitigation", "Identified risks, impact levels and contingency plans", "11")

flat_rect(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.45), PRIMARY)
for x, w, h in [(Inches(0.9), Inches(0.5), "No."), (Inches(1.4), Inches(2.8), "Risk Description"),
                 (Inches(4.3), Inches(1.0), "Likelihood"), (Inches(5.4), Inches(1.0), "Impact"),
                 (Inches(6.5), Inches(5.5), "Mitigation Strategy")]:
    txt(slide, x, Inches(1.68), w, Inches(0.4), h, size=11, color=WHITE, bold=True)

risks = [
    ("R1", "Database connectivity failure in production", "Medium", "High",
     "Graceful fallback to development mode. Connection retry with exponential backoff. Health check endpoints for monitoring."),
    ("R2", "Brute force attack on login endpoint", "High", "High",
     "Rate limiting (10/min/IP server-side, 5 attempts client lockout). bcrypt hashing with salt. Full login attempt audit logging."),
    ("R3", "Data loss or corruption", "Low", "Critical",
     "Supabase managed PostgreSQL with automated backups. MongoDB Atlas point-in-time recovery. Versioned migration scripts."),
    ("R4", "USSD service unavailability", "Medium", "Medium",
     "Web dashboard and mobile app remain functional. USSD interaction logs preserved for replay. Independent service architecture."),
    ("R5", "Unauthorized access to sensitive data", "Low", "Critical",
     "Role-based access control. Input sanitization against XSS/injection. Security headers (CSP, X-Frame-Options). Encrypted sessions."),
    ("R6", "Staff turnover or capacity constraints", "Medium", "Medium",
     "Maximum 3-user model simplifies onboarding. Self-service signup for PA/Clerk. AI assistant provides in-dashboard guidance."),
    ("R7", "Scalability under increased usage", "Low", "Medium",
     "PostgreSQL connection pool (20 connections). Redis caching layer. Containerized deployment with horizontal scaling capability."),
]

y = Inches(2.15)
for i, (num, desc, likelihood, impact, mitigation) in enumerate(risks):
    bg = SECTION_BG if i % 2 == 0 else BG_WHITE
    flat_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), bg)
    txt(slide, Inches(0.9), y + Pt(3), Inches(0.5), Inches(0.3), num, size=10, color=ACCENT, bold=True, font='Consolas')
    txt(slide, Inches(1.4), y + Pt(3), Inches(2.8), Inches(0.6), desc, size=10, color=DARK_GRAY)
    l_color = AMBER_WARN if likelihood == "Medium" else (RED_RISK if likelihood == "High" else GREEN_OK)
    txt(slide, Inches(4.3), y + Pt(3), Inches(1.0), Inches(0.3), likelihood, size=10, color=l_color, bold=True)
    i_color = RED_RISK if impact in ("High", "Critical") else AMBER_WARN
    txt(slide, Inches(5.4), y + Pt(3), Inches(1.0), Inches(0.3), impact, size=10, color=i_color, bold=True)
    txt(slide, Inches(6.5), y + Pt(3), Inches(5.5), Inches(0.6), mitigation, size=9, color=DARK_GRAY, spacing=1.3)
    y += Inches(0.7)
footer(slide)


# =====================================================================
# SLIDE 15 — CONCLUSION AND RECOMMENDATIONS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
header_bar(slide, "Conclusion and Recommendations", "Summary and proposed next steps", "12")

txt(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.8),
    "The VOO Ward Admin Dashboard provides Kyamatu Ward with a comprehensive, secure, and scalable digital governance platform. By centralizing issue management, bursary processing, citizen communication, and analytics into a single system accessible via web, USSD, and mobile channels, the platform directly addresses the identified challenges of fragmented processes, limited transparency, and data-driven decision making.",
    size=13, color=DARK_GRAY, spacing=1.6)

card(slide, Inches(0.8), Inches(2.8), Inches(5.7), Inches(2.0),
     "Key Deliverables", [
         "- Web-based Admin Dashboard with 6 core modules",
         "- USSD service integration for feature phone access",
         "- Role-based access control (MCA, PA, Clerk)",
         "- Real-time analytics with charts and geographic mapping",
         "- Comprehensive audit trail and CSV export capabilities",
         "- AI chat assistant for in-dashboard support",
     ], PRIMARY)

card(slide, Inches(6.7), Inches(2.8), Inches(5.8), Inches(2.0),
     "Technical Highlights", [
         "- Node.js/Express.js backend with PostgreSQL and MongoDB",
         "- Docker containerization for consistent deployment",
         "- Enterprise-grade security (bcrypt, CSP, rate limiting)",
         "- Responsive design supporting desktop, tablet, and mobile",
         "- Graceful degradation when external services are unavailable",
         "- Real-time updates with 30-second polling and notifications",
     ], ACCENT)

txt(slide, Inches(0.8), Inches(5.1), Inches(5), Inches(0.3),
    "Recommended Next Steps", size=16, color=PRIMARY, bold=True)
flat_rect(slide, Inches(0.8), Inches(5.42), Inches(2), Pt(3), ACCENT)

steps = [
    ("1.", "Stakeholder Review", "Present this proposal to ward leadership for feedback and approval."),
    ("2.", "Requirements Sign-Off", "Confirm functional requirements and any additional customizations needed."),
    ("3.", "Development Kickoff", "Begin Phase 1 (Foundation) with database setup and authentication framework."),
    ("4.", "Iterative Delivery", "Deploy incremental releases at the end of each phase for stakeholder review."),
]
y = Inches(5.65)
for num, title, desc in steps:
    rect(slide, Inches(0.8), y, Inches(0.35), Inches(0.35), ACCENT)
    txt(slide, Inches(0.82), y + Pt(1), Inches(0.33), Inches(0.33),
        num, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, Inches(1.3), y, Inches(2), Inches(0.28),
        title, size=12, color=PRIMARY, bold=True)
    txt(slide, Inches(3.5), y + Pt(1), Inches(8.5), Inches(0.28),
        desc, size=11, color=DARK_GRAY)
    y += Inches(0.38)
footer(slide)


# =====================================================================
# SLIDE 16 — THANK YOU / Q&A
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, PRIMARY)
flat_rect(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5), ACCENT)

txt(slide, Inches(2), Inches(1.5), Inches(9.333), Inches(0.5),
    "PROJECT PROPOSAL", size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(1.0),
    "Thank You", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
flat_rect(slide, Inches(5.5), Inches(3.4), Inches(2.333), Pt(3), ACCENT)
txt(slide, Inches(2), Inches(3.7), Inches(9.333), Inches(0.5),
    "VOO Ward Admin Dashboard", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(2), Inches(4.3), Inches(9.333), Inches(0.5),
    "Comprehensive Digital Governance Platform for Kyamatu Ward",
    size=14, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)
flat_rect(slide, Inches(4), Inches(5.2), Inches(5.333), Pt(1), RGBColor(0x3A, 0x5A, 0x7A))
txt(slide, Inches(2), Inches(5.5), Inches(9.333), Inches(0.5),
    "Questions and Discussion", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(2), Inches(6.2), Inches(9.333), Inches(0.5),
    "VOO Ward Development Team  |  February 2026  |  Confidential",
    size=11, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)


# =====================================================================
# SAVE
# =====================================================================
output_path = os.path.join(os.path.dirname(__file__), "docs", "VOO_Ward_Admin_Dashboard_Proposal.pptx")
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
