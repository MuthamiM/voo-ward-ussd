#!/usr/bin/env python3
"""Convert PPTX files into detailed, professional DOCX files."""

import os
import re
from datetime import datetime

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from pptx import Presentation


FOOTER_NOISE = {
    "VOO WARD",
    "Confidential",
    "VOO Ward Admin Dashboard",
}

BRAND_PRIMARY = RGBColor(31, 73, 125)
BRAND_ACCENT = RGBColor(54, 96, 146)


def apply_brand_styles(doc):
    """Apply consistent branded styling to the Word document."""
    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Calibri"
    normal_style.font.size = Pt(11)

    h1 = doc.styles["Heading 1"]
    h1.font.name = "Calibri"
    h1.font.size = Pt(16)
    h1.font.bold = True
    h1.font.color.rgb = BRAND_PRIMARY

    h2 = doc.styles["Heading 2"]
    h2.font.name = "Calibri"
    h2.font.size = Pt(13)
    h2.font.bold = True
    h2.font.color.rgb = BRAND_ACCENT

    h3 = doc.styles["Heading 3"]
    h3.font.name = "Calibri"
    h3.font.size = Pt(11)
    h3.font.bold = True
    h3.font.color.rgb = BRAND_ACCENT


def add_section_divider(doc):
    """Add a subtle section divider line."""
    divider = doc.add_paragraph("-" * 60)
    divider.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if divider.runs:
        divider.runs[0].font.color.rgb = BRAND_ACCENT
        divider.runs[0].font.size = Pt(9)


def append_field_run(paragraph, field_code, placeholder=""):
    """Insert a dynamic Word field into a paragraph."""
    run = paragraph.add_run()
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")

    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = field_code

    fld_separate = OxmlElement("w:fldChar")
    fld_separate.set(qn("w:fldCharType"), "separate")

    run._r.append(fld_begin)
    run._r.append(instr)
    run._r.append(fld_separate)

    if placeholder:
        paragraph.add_run(placeholder)

    fld_end_run = paragraph.add_run()
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    fld_end_run._r.append(fld_end)


def add_toc(doc):
    """Add an auto-updating table of contents field."""
    toc_title = doc.add_heading("Table of Contents", level=1)
    toc_title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    toc_para = doc.add_paragraph()
    append_field_run(toc_para, r'TOC \\o "1-3" \\h \\z \\u', "Update field to generate TOC.")
    doc.add_page_break()


def configure_header_footer(doc, doc_title):
    """Apply a formal header and footer with dynamic page numbering."""
    section = doc.sections[0]

    header_para = section.header.paragraphs[0] if section.header.paragraphs else section.header.add_paragraph()
    header_para.text = doc_title
    header_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    if header_para.runs:
        header_para.runs[0].font.size = Pt(9)
        header_para.runs[0].font.color.rgb = BRAND_ACCENT

    footer_para = section.footer.paragraphs[0] if section.footer.paragraphs else section.footer.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_para.add_run("Page ")
    append_field_run(footer_para, "PAGE")
    footer_para.add_run(" of ")
    append_field_run(footer_para, "NUMPAGES")
    for run in footer_para.runs:
        run.font.size = Pt(9)


def sanitize_text(text):
    """Remove control chars and emoji-like symbols while keeping readable text."""
    if not text:
        return ""

    cleaned = []
    for ch in text:
        code = ord(ch)
        if ch in "\n\r\t":
            cleaned.append(ch)
            continue
        if code < 32:
            continue
        # Remove common emoji/symbol blocks to keep professional tone.
        if (
            0x1F300 <= code <= 0x1FAFF
            or 0x2600 <= code <= 0x27BF
            or 0xFE00 <= code <= 0xFE0F
        ):
            continue
        cleaned.append(ch)

    text = "".join(cleaned)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def shape_sort_key(shape):
    """Sort shapes in visual reading order: top-to-bottom, left-to-right."""
    return (getattr(shape, "top", 0), getattr(shape, "left", 0))


def extract_slide_lines(slide):
    """Extract slide content as ordered lines with level and size hints."""
    lines = []
    for shape in sorted(slide.shapes, key=shape_sort_key):
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                raw = paragraph.text or ""
                text = sanitize_text(raw)
                if not text:
                    continue

                if text in FOOTER_NOISE or text.upper() in FOOTER_NOISE:
                    continue

                font_size = None
                for run in paragraph.runs:
                    if run.font and run.font.size:
                        font_size = run.font.size.pt
                        break

                lines.append(
                    {
                        "text": text,
                        "level": paragraph.level if paragraph.level is not None else 0,
                        "font_size": font_size,
                    }
                )

        if getattr(shape, "has_table", False):
            table_rows = []
            table = shape.table
            for row in table.rows:
                row_cells = [sanitize_text(cell.text) for cell in row.cells]
                table_rows.append(row_cells)
            if table_rows:
                lines.append({"table": table_rows})

    return lines


def add_table(doc, rows):
    """Add table content to a Word document."""
    if not rows:
        return
    col_count = max(len(r) for r in rows) if rows else 0
    if col_count == 0:
        return

    table = doc.add_table(rows=len(rows), cols=col_count)
    table.style = "Table Grid"

    for r_idx, row in enumerate(rows):
        for c_idx in range(col_count):
            value = row[c_idx] if c_idx < len(row) else ""
            table.cell(r_idx, c_idx).text = value


def add_slide_content(doc, slide_number, lines):
    """Render one slide as a structured section in DOCX."""
    if not lines:
        return

    # Use first non-empty text line as slide title when possible.
    title_text = None
    for item in lines:
        if "text" in item:
            title_text = item["text"]
            break

    heading = f"Slide {slide_number}"
    if title_text:
        heading = f"Slide {slide_number}: {title_text}"
    head = doc.add_heading(heading, level=1)
    head.alignment = WD_ALIGN_PARAGRAPH.LEFT
    add_section_divider(doc)

    skipped_title = False
    for item in lines:
        if "table" in item:
            add_table(doc, item["table"])
            doc.add_paragraph("")
            continue

        text = item["text"]
        level = item["level"]
        size = item["font_size"]

        if not skipped_title and title_text and text == title_text:
            skipped_title = True
            continue

        # Secondary heading detection by size and brevity.
        if size and size >= 18 and len(text) <= 90:
            doc.add_heading(text, level=2)
            continue

        # Numbered lines become numbered list entries.
        if re.match(r"^\d+[\.)\-:]\s+", text):
            doc.add_paragraph(text, style="List Number")
            continue

        # Bullet-like lines or indented paragraphs become bullets.
        if text.startswith(("- ", "* ", "• ")) or level > 0:
            clean = re.sub(r"^[-*•]\s*", "", text)
            doc.add_paragraph(clean, style="List Bullet")
            continue

        doc.add_paragraph(text)

    doc.add_page_break()


def create_professional_docx(pptx_file, output_file):
    """Create a detailed DOCX from a PPTX file."""
    try:
        prs = Presentation(pptx_file)
    except Exception as exc:
        print(f"Error reading {pptx_file}: {exc}")
        return False

    doc = Document()
    apply_brand_styles(doc)

    file_title = os.path.splitext(os.path.basename(pptx_file))[0].replace("_", " ")
    configure_header_footer(doc, file_title)

    cover = doc.add_heading(file_title, level=0)
    cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if cover.runs:
        cover.runs[0].font.color.rgb = BRAND_PRIMARY
    sub = doc.add_paragraph("Detailed Document Version")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if sub.runs:
        sub.runs[0].font.color.rgb = BRAND_ACCENT
        sub.runs[0].font.size = Pt(12)
    gen = doc.add_paragraph(f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    gen.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if gen.runs:
        gen.runs[0].font.size = Pt(10)
    add_section_divider(doc)
    doc.add_page_break()

    add_toc(doc)

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            lines = extract_slide_lines(slide)
            add_slide_content(doc, idx, lines)
        except Exception as exc:
            doc.add_heading(f"Slide {idx}", level=1)
            doc.add_paragraph(f"Slide content could not be fully processed: {exc}")
            doc.add_page_break()

    try:
        doc.save(output_file)
        return True
    except Exception as exc:
        print(f"Error saving {output_file}: {exc}")
        return False


if __name__ == "__main__":
    docs_dir = r"c:\Users\kivuv\Documents\voo-ward-ussd\docs"
    pptx_files = [
        "VOO_Ward_Admin_Dashboard_Proposal.pptx",
        "VOO_Ward_Admin_Dashboard.pptx",
        "VOO_Ward_Admin_Dashboard_v2.pptx",
    ]

    for pptx_name in pptx_files:
        pptx_path = os.path.join(docs_dir, pptx_name)
        if not os.path.exists(pptx_path):
            print(f"Missing: {pptx_name}")
            continue

        output_name = pptx_name.replace(".pptx", "_professional.docx")
        output_path = os.path.join(docs_dir, output_name)
        if create_professional_docx(pptx_path, output_path):
            print(f"Created: {output_name}")
        else:
            print(f"Failed: {pptx_name}")
