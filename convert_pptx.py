#!/usr/bin/env python3
"""
Convert PowerPoint presentations to Word documents
"""

from pptx import Presentation
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import os

def clean_text(text):
    """Clean text to remove null bytes and control characters"""
    if not text:
        return ""
    # Remove null bytes and control characters
    return ''.join(char for char in text if ord(char) >= 32 or char in '\n\r\t')

def convert_pptx_to_docx(pptx_file, output_file):
    """Convert PowerPoint presentation to Word document"""
    
    # Load presentation
    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        print(f"Error reading {pptx_file}: {e}")
        return False
    
    # Create document
    doc = Document()
    
    # Add title from filename
    title = os.path.splitext(os.path.basename(pptx_file))[0]
    title = title.replace('_', ' ')
    title = clean_text(title)
    doc.add_heading(title, level=0)
    doc.add_paragraph()
    
    # Process each slide
    for slide_num, slide in enumerate(prs.slides, 1):
        doc.add_heading(f'Slide {slide_num}', level=1)
        
        # Extract text from all shapes in the slide
        text_found = False
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = clean_text(shape.text.strip())
                
                # Skip slide titles if they're the same as the main title or empty
                if text and len(text) > 0:
                    # Check if it's likely a heading (shorter text in larger font)
                    if hasattr(shape, "text_frame"):
                        frame = shape.text_frame
                        if frame.paragraphs:
                            para = frame.paragraphs[0]
                            if para.runs and len(text) < 100:
                                # Likely a title/heading
                                doc.add_heading(text, level=2)
                                text_found = True
                            else:
                                # Regular content
                                p = doc.add_paragraph(text)
                                text_found = True
                    else:
                        p = doc.add_paragraph(text)
                        text_found = True
            
            # Extract tables
            if shape.has_table:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                
                doc_table = doc.add_table(rows=rows, cols=cols)
                doc_table.style = 'Light Grid Accent 1'
                
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        doc_cell = doc_table.rows[row_idx].cells[col_idx]
                        doc_cell.text = clean_text(cell.text)
                
                text_found = True
        
        if not text_found:
            doc.add_paragraph("[No text content on this slide]")
        
        doc.add_paragraph()  # Spacing between slides
    
    # Save document
    doc.save(output_file)
    return True

if __name__ == "__main__":
    docs_dir = r"c:\Users\kivuv\Documents\voo-ward-ussd\docs"
    
    pptx_files = [
        "VOO_Ward_Admin_Dashboard.pptx",
        "VOO_Ward_Admin_Dashboard_Proposal.pptx",
        "VOO_Ward_Admin_Dashboard_v2.pptx"
    ]
    
    for pptx in pptx_files:
        pptx_path = os.path.join(docs_dir, pptx)
        if os.path.exists(pptx_path):
            output_name = pptx.replace('.pptx', '.docx')
            output_path = os.path.join(docs_dir, output_name)
            
            if convert_pptx_to_docx(pptx_path, output_path):
                print(f"✓ Converted: {pptx} → {output_name}")
            else:
                print(f"✗ Failed: {pptx}")
        else:
            print(f"✗ Not found: {pptx_path}")
